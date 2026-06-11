# -*- coding: utf-8 -*-
"""Build the VendorPulse Solution Architecture slide (Shell).

Monochrome engineering style (matches the client's human-made reference deck):
white boxes, thin grey borders, Calibri, grey section headers, a single red
accent reserved for the numbered request-flow markers and the human-approval
control.  Client -> Edge -> (Shell Azure Subscription: App Service compute tier
+ Data tier) -> External services.  Icon placeholders are dashed squares the
author fills in manually.   Run:  python _build_arch_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# --- monochrome palette + single red accent ---
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x20, 0x20, 0x20)   # box label text (near-black)
GREY   = RGBColor(0x59, 0x59, 0x59)   # section headers / sub text
BORDER = RGBColor(0x40, 0x40, 0x40)   # thin box outlines
HAIR   = RGBColor(0xBF, 0xBF, 0xBF)   # light rules / icon placeholders
RED    = RGBColor(0xC0, 0x00, 0x00)   # the only accent: flow numbers + approval control
FONT   = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def shape(kind, l, t, w, h, fill, line=BORDER, width=0.75, dash=False):
    sh = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(width)
        if dash:
            ln = sh._element.spPr.get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sh.shadow.inherit = False
    return sh


def rect(l, t, w, h, fill=WHITE, line=BORDER, width=0.75, dash=False):
    return shape(MSO_SHAPE.RECTANGLE, l, t, w, h, fill, line, width, dash)


def write(sh, title, sub=None, tcolor=INK, tsize=10, ssize=8, align=PP_ALIGN.CENTER, bold=True):
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = None
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = title; r.font.bold = bold; r.font.size = Pt(tsize); r.font.name = FONT; r.font.color.rgb = tcolor
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = align
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(ssize); r2.font.name = FONT; r2.font.color.rgb = GREY


def textbox(l, t, w, h, text, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = FONT; r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return tb


def connect(x1, y1, x2, y2, color=BORDER, width=1.0, arrow=True, dash=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def numcircle(cx, cy, n, d=0.3):
    """Red numbered flow marker (the single accent colour)."""
    o = shape(MSO_SHAPE.OVAL, cx, cy, d, d, RED, line=None)
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.bold = True; r.font.size = Pt(11); r.font.name = FONT; r.font.color.rgb = WHITE
    return o


def icon_ph(l, t, d=0.32):
    """Dashed placeholder square for an icon the author drops in manually."""
    sq = rect(l, t, d, d, WHITE, line=HAIR, dash=True)
    write(sq, "icon", tcolor=HAIR, tsize=6, bold=False)
    return sq


def lanehead(l, t, w, text, num=None):
    """Grey bold section header with a thin underline rule (human reference style)."""
    textbox(l, t, w, 0.28, text, size=12, color=GREY, bold=True)
    connect(l, t + 0.34, l + w, t + 0.34, color=HAIR, width=1.0, arrow=False)
    if num is not None:
        numcircle(l - 0.02, t - 0.30, num)


# ===================== canvas =====================
rect(0, 0, 13.333, 7.5, WHITE, line=None)
textbox(0.4, 0.18, 11.5, 0.5, "VendorPulse Solution Architecture (Shell)", size=22, color=INK, bold=True)
textbox(0.4, 0.66, 12.5, 0.3,
        "Single-tenant in Shell Azure.  Deterministic core first, AI second.  Agent layer: MAF SDK on Microsoft Foundry.",
        size=10.5, color=GREY)
connect(0.4, 1.02, 12.93, 1.02, color=HAIR, width=1.0, arrow=False)

# ---------- CLIENT lane ----------
lanehead(0.35, 1.30, 2.4, "Client", num=1)
icon_ph(0.40, 1.72)
cl = rect(0.35, 2.20, 2.4, 0.42); write(cl, "Web client (browser)", tcolor=INK, tsize=10.5)
for (t, sub, y, h) in [
    ("VMO Coordinator / Sponsor / Viewer", "authenticated via Entra SSO", 2.80, 0.85),
    ("React 19 SPA", "Design System / ApprovalPanel", 3.85, 0.85),
    ("Native scorecard form", "magic-link, Entra SSO (ADR-005)", 4.90, 0.85),
    ("AI-generated, pending approval", "transparency badge (IRM 3.5.3)", 5.95, 0.85),
]:
    redline = (t.startswith("AI-generated"))
    c = rect(0.35, y, 2.4, h, WHITE, line=RED if redline else BORDER, width=1.25 if redline else 0.75)
    write(c, t, sub, tcolor=INK)

# ---------- EDGE ----------
lanehead(2.95, 1.30, 1.10, "Edge")
icon_ph(3.32, 1.72)
e = rect(2.95, 3.20, 1.05, 1.40); write(e, "Azure Front Door + WAF", "TLS, OWASP, origin-lock", tcolor=INK, tsize=9.5, ssize=7)

# ---------- SHELL AZURE SUBSCRIPTION boundary ----------
AZx, AZy, AZw, AZh = 4.25, 1.92, 5.75, 5.13
lanehead(AZx, 1.30, AZw, "Shell Azure Subscription", num=3)
rect(AZx, AZy, AZw, AZh, WHITE, line=BORDER, width=1.25)
textbox(AZx + 0.12, AZy + 0.04, AZw - 0.24, 0.26, "West Europe  /  Private VNet", size=8.5, color=GREY, bold=True)

# --- App Service (compute tier) ---
ASx, ASy, ASw, ASh = AZx + 0.2, AZy + 0.45, AZw - 0.4, 2.30
rect(ASx, ASy, ASw, ASh, WHITE, line=BORDER)
icon_ph(ASx + 0.06, ASy + 0.05, 0.26)
textbox(ASx + 0.36, ASy + 0.04, ASw - 0.4, 0.30, "Azure App Service (Linux container) - FastAPI app", size=9.5, color=GREY, bold=True)
mods = [
    ("Entra OIDC auth + RBAC", False),
    ("WorkflowEngine (12-state)", False),
    ("Approval gate", True),          # red border: the human-approval control
    ("MAF Agent layer to Foundry", False),
    ("Deterministic services", False),
    ("GraphService (app-only cert)", False),
]
gx0, gy0 = ASx + 0.12, ASy + 0.42; cw = (ASw - 0.36) / 2; chh = 0.52; gap = 0.08
for i, (t, redline) in enumerate(mods):
    col = i % 2; row = i // 2
    c = rect(gx0 + col * (cw + 0.12), gy0 + row * (chh + gap), cw, chh, WHITE,
             line=RED if redline else BORDER, width=1.25 if redline else 0.75)
    write(c, t, tcolor=INK, tsize=9)

# --- Data tier ---
DTx, DTy, DTw, DTh = AZx + 0.2, ASy + ASh + 0.15, AZw - 0.4, 1.78
rect(DTx, DTy, DTw, DTh, WHITE, line=BORDER)
textbox(DTx + 0.08, DTy + 0.04, DTw - 0.16, 0.26, "Data Tier - Managed Identity, Private Link (no public access)", size=9, color=GREY, bold=True)
db = shape(MSO_SHAPE.CAN, DTx + 0.18, DTy + 0.42, 1.5, 1.10, WHITE, line=BORDER)
write(db, "Azure PostgreSQL", "Flexible Server", tcolor=INK, tsize=9.5, ssize=7)
textbox(DTx + 0.1, DTy + 1.52, 1.7, 0.2, "cycles / scorecards / meetings / agent_runs", size=6.5, color=GREY, align=PP_ALIGN.CENTER)
others = [("Azure Key Vault", "LLM key / Graph cert / JWT"),
          ("Blob Storage", "minutes, transcripts"),
          ("App Insights + Log Analytics", "OTel / immutable audit")]
ox = DTx + 1.85; ow = (DTw - 1.85 - 0.12) / 3 - 0.1
for i, (t, sub) in enumerate(others):
    c = rect(ox + i * (ow + 0.12), DTy + 0.42, ow, 1.10, WHITE, line=BORDER); write(c, t, sub, tcolor=INK, tsize=8.5, ssize=6.8)

# ---------- EXTERNAL services lane ----------
EXx, EXy, EXw, EXh = 10.15, 1.92, 2.83, 5.13
lanehead(EXx, 1.30, EXw, "External", num=5)
rect(EXx, EXy, EXw, EXh, WHITE, line=BORDER, width=1.25)
textbox(EXx + 0.12, EXy + 0.04, EXw - 0.24, 0.26, "Outbound HTTPS", size=8.5, color=GREY, bold=True)
for (t, sub, y, redline) in [
    ("Microsoft Entra ID", "Shell SSO, groups to roles, app identity", 2.45, False),
    ("Microsoft Foundry", "Responses API + content safety", 3.85, False),
    ("Microsoft Graph (Shell tenant)", "Mail.Send / Calendars / OnlineMtgs", 5.25, False),
]:
    c = rect(EXx + 0.15, y, EXw - 0.3, 1.15, WHITE, line=BORDER)
    icon_ph(EXx + 0.24, y + 0.10, 0.26)
    write(c, t, sub, tcolor=INK, tsize=9.5)

# ---------- flow connectors (thin, monochrome) ----------
connect(2.77, 3.90, 2.93, 3.90)                 # client -> edge
connect(4.02, 3.90, 4.23, 3.90)                 # edge -> azure
connect(AZx + AZw / 2, ASy + ASh, AZx + AZw / 2, DTy)  # compute -> data (down)
connect(10.00, 3.00, 10.13, 3.00, color=BORDER)        # app -> external
textbox(9.45, 2.62, 1.7, 0.2, "outbound HTTPS", size=6.5, color=GREY, align=PP_ALIGN.CENTER)
numcircle(2.78, 3.30, 2)        # edge marker
numcircle(AZx + AZw / 2 - 0.15, DTy - 0.34, 4)  # data marker

# ---------- legend (numbered request flow) ----------
LGx, LGy = 0.4, 7.06
legend = ("1 Client / SSO sign-in    2 Front Door + WAF    3 App Service: workflow + approval gate"
          "    4 Data tier via Private Link    5 External services over outbound HTTPS")
textbox(LGx, LGy, 12.5, 0.3, legend, size=8, color=GREY, italic=False)

OUT = "VendorPulse_Solution_Architecture.pptx"; prs.save(OUT); print("wrote", OUT)
