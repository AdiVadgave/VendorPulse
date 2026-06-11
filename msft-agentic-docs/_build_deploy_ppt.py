# -*- coding: utf-8 -*-
"""Build the VendorPulse Deployment Architecture slide (Shell Azure).

Monochrome engineering style matching the client's human-made reference deck:
white boxes, thin grey borders, Calibri, grey section headers, dashed cloud/VNet
boundaries, and a single red accent reserved for the numbered deployment-flow
markers and the production app-identity control.  Icon placeholders are dashed
squares the author fills in manually.   Run: python _build_deploy_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x20, 0x20, 0x20)
GREY   = RGBColor(0x59, 0x59, 0x59)
BORDER = RGBColor(0x40, 0x40, 0x40)
HAIR   = RGBColor(0xBF, 0xBF, 0xBF)
RED    = RGBColor(0xC0, 0x00, 0x00)
FONT   = "Calibri"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])


def shape(kind, l, t, w, h, fill, line=BORDER, width=0.75, dash=False):
    sh = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(width)
        if dash:
            ln = sh._element.spPr.get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    sh.shadow.inherit = False
    return sh


def rect(l, t, w, h, fill=WHITE, line=BORDER, width=0.75, dash=False):
    return shape(MSO_SHAPE.RECTANGLE, l, t, w, h, fill, line, width, dash)


def write(sh, title, sub=None, tcolor=INK, tsize=9.5, ssize=7, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True):
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor; tf.auto_size = None
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = title; r.font.bold = bold; r.font.size = Pt(tsize); r.font.name = FONT; r.font.color.rgb = tcolor
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = align
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(ssize); r2.font.name = FONT; r2.font.color.rgb = GREY


def textbox(l, t, w, h, text, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.name = FONT; r.font.bold = bold; r.font.color.rgb = color; r.font.italic = italic
    return tb


def connect(x1, y1, x2, y2, color=BORDER, width=1.0, arrow=True, dash=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width); cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def numcircle(cx, cy, n, d=0.3):
    o = shape(MSO_SHAPE.OVAL, cx, cy, d, d, RED, line=None)
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.bold = True; r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = WHITE
    return o


def icon_ph(l, t, d=0.30):
    sq = rect(l, t, d, d, WHITE, line=HAIR, dash=True)
    write(sq, "icon", tcolor=HAIR, tsize=6, bold=False)
    return sq


def hdr(l, t, w, text, h=0.32, tsize=8.8):
    """White header band with thin border + grey label (no colour fill)."""
    b = rect(l, t, w, h, WHITE, line=BORDER); write(b, text, tcolor=GREY, tsize=tsize); return b


# ===================== canvas =====================
rect(0, 0, 13.333, 7.5, WHITE, line=None)
textbox(0.4, 0.18, 9.8, 0.45, "VendorPulse Deployment Architecture (Shell Azure)", size=20, color=INK, bold=True)
textbox(0.4, 0.62, 10.0, 0.3, "Single-tenant.  West Europe.  Azure App Service.  Private VNet.  App-only certificate auth.", size=10, color=GREY)
connect(0.4, 0.98, 12.93, 0.98, color=HAIR, width=1.0, arrow=False)
lg = rect(11.5, 0.20, 1.5, 0.60, WHITE, line=HAIR, dash=True); write(lg, "[ PRODUCT LOGO ]", tcolor=HAIR, tsize=8)

# ---------- LEFT: user + SSO path ----------
textbox(0.35, 1.18, 1.95, 0.24, "Access path", size=11, color=GREY, bold=True)
connect(0.35, 1.46, 2.30, 1.46, color=HAIR, width=1.0, arrow=False)


def leftcard(t, sub, y, h, icon=False):
    if icon:
        icon_ph(0.40, y + (h - 0.26) / 2, 0.26)
    c = rect(0.35, y, 1.95, h, WHITE, line=BORDER); write(c, t, sub, tcolor=INK, tsize=8.8, ssize=6.8); return c


leftcard("Shell user", "on corporate network", 1.55, 0.7, icon=True)
leftcard("Microsoft Entra ID", "Shell SSO", 2.40, 0.62)
leftcard("Single Sign-On", "OIDC / MSAL", 3.12, 0.62)
leftcard("vendorpulse.it.shell.com", "DNS to Front Door", 3.84, 0.62)
leftcard("Shell egress proxy", "outbound control", 4.90, 0.62)
numcircle(0.05, 1.62, 1)

# ---------- AZURE CLOUD boundary (dashed) ----------
AZx, AZy, AZw, AZh = 2.5, 1.4, 7.95, 5.55
rect(AZx, AZy, AZw, AZh, WHITE, line=BORDER, width=1.0, dash=True)
icon_ph(AZx + 0.12, AZy + 0.06, 0.24)
textbox(AZx + 0.42, AZy + 0.06, 5.0, 0.26, "Shell Azure Subscription - West Europe", size=9, color=GREY, bold=True)
numcircle(AZx - 0.02, AZy + 0.30, 2)
# Front Door band
hdr(AZx + 0.15, AZy + 0.40, AZw - 0.3, "Azure Front Door + WAF   -   TLS 1.2+, OWASP, origin-locked to App Service", h=0.32, tsize=8.8)

# VNet box (dashed)
VNx, VNy, VNw, VNh = AZx + 0.15, AZy + 0.82, AZw - 0.3, 2.7
rect(VNx, VNy, VNw, VNh, WHITE, line=BORDER, width=1.0, dash=True)
icon_ph(VNx + 0.08, VNy + 0.05, 0.22)
textbox(VNx + 0.36, VNy + 0.05, 4.0, 0.24, "Private VNet (VNet integration, Private Endpoints)", size=8.2, color=GREY, bold=True)
numcircle(VNx - 0.02, VNy + 0.02, 3)

# App Service box
ASx, ASy, ASw, ASh = VNx + 0.18, VNy + 0.34, 4.25, 2.20
rect(ASx, ASy, ASw, ASh, WHITE, line=BORDER)
hdr(ASx, ASy, ASw, "Azure App Service (Linux container)", h=0.32, tsize=8.8)
ps = rect(ASx + 0.12, ASy + 0.40, ASw - 0.24, 1.05, WHITE, line=BORDER)
write(ps, "Production slot  (TLS 1.2+)", tcolor=GREY, tsize=8, anchor=MSO_ANCHOR.TOP)
r1 = rect(ASx + 0.25, ASy + 0.70, ASw - 0.5, 0.32, WHITE, line=BORDER); write(r1, "SPA (React / nginx)", tcolor=INK, tsize=7.5)
r2 = rect(ASx + 0.25, ASy + 1.06, ASw - 0.5, 0.32, WHITE, line=BORDER); write(r2, "FastAPI + MAF agents", tcolor=INK, tsize=7.5)
ss = rect(ASx + 0.12, ASy + 1.52, ASw - 0.24, 0.55, WHITE, line=BORDER)
write(ss, "Staging slot", "blue-green / zero-downtime swap", tcolor=GREY, tsize=7.5, ssize=6.5)

# Data RSG box
DRx, DRy, DRw, DRh = ASx + ASw + 0.35, VNy + 0.34, VNw - (ASx + ASw + 0.35 - VNx) - 0.18, 2.20
rect(DRx, DRy, DRw, DRh, WHITE, line=BORDER)
hdr(DRx, DRy, DRw, "Data Resource Group", h=0.32, tsize=8.8)
numcircle(DRx - 0.16, DRy + 0.01, 4)
db = shape(MSO_SHAPE.CAN, DRx + 0.2, DRy + 0.44, 0.95, 1.05, WHITE, line=BORDER)
write(db, "Postgres", "Flexible", tcolor=INK, tsize=7.5, ssize=6)
textbox(DRx + 0.05, DRy + 1.52, 1.25, 0.18, "Private Link", size=6.2, color=GREY, align=PP_ALIGN.CENTER)
kv = rect(DRx + 1.3, DRy + 0.44, DRw - 1.45, 0.48, WHITE, line=BORDER); write(kv, "Key Vault", "cert, keys (PE)", tcolor=INK, tsize=7.5, ssize=6)
bl = rect(DRx + 1.3, DRy + 0.98, DRw - 1.45, 0.5, WHITE, line=BORDER); write(bl, "Blob Storage", "minutes (PE)", tcolor=INK, tsize=7.5, ssize=6)

# Ops / CI-CD icon row
ORy = VNy + VNh + 0.14
textbox(AZx + 0.15, ORy - 0.02, 3.0, 0.2, "DevOps & Observability", size=8, color=GREY, bold=True)
connect(AZx + 0.15, ORy + 0.18, AZx + AZw - 0.15, ORy + 0.18, color=HAIR, width=1.0, arrow=False)
ops = ["GitHub Actions / Azure DevOps", "Container Registry (ACR)", "Azure Monitor", "Application Insights", "Log Analytics / immutable audit"]
ow = (AZw - 0.3 - 0.4) / 5; ox = AZx + 0.15
for i, t in enumerate(ops):
    c = rect(ox + i * (ow + 0.1), ORy + 0.24, ow, 0.66, WHITE, line=BORDER); write(c, t, tcolor=INK, tsize=7.5)

# ---------- RIGHT: external services ----------
EXx, EXy, EXw = 10.62, 1.4, 2.4
rect(EXx, EXy, EXw, 3.0, WHITE, line=BORDER, width=1.0, dash=True)
textbox(EXx + 0.1, EXy + 0.06, 2.2, 0.24, "External (outbound HTTPS)", size=8.2, color=GREY, bold=True)
numcircle(EXx - 0.02, EXy + 0.30, 5)
ec1 = rect(EXx + 0.15, EXy + 0.40, EXw - 0.3, 0.78, WHITE, line=BORDER); icon_ph(EXx + 0.24, EXy + 0.50, 0.22); write(ec1, "Microsoft Foundry", "Responses API + content safety", tcolor=INK, tsize=8.5, ssize=6.8)
ec2 = rect(EXx + 0.15, EXy + 1.26, EXw - 0.3, 0.78, WHITE, line=BORDER); icon_ph(EXx + 0.24, EXy + 1.36, 0.22); write(ec2, "Microsoft Graph (Shell)", "Mail / Calendar / Teams (app-only cert)", tcolor=INK, tsize=8.5, ssize=6.8)
ec3 = rect(EXx + 0.15, EXy + 2.12, EXw - 0.3, 0.78, WHITE, line=RED, width=1.25); write(ec3, "Entra ID (app reg)", "VendorPulse-Prod, cert in KV", tcolor=INK, tsize=8.5, ssize=6.8)

# ---------- flow connectors ----------
connect(2.32, 2.12, 2.50, 2.12, color=BORDER)                       # left -> azure (Front Door)
connect(ASx + ASw, ASy + 0.9, DRx, ASy + 0.9, color=BORDER)         # app -> data
textbox(ASx + ASw - 0.05, ASy + 0.58, 0.6, 0.2, "Private Link", size=5.8, color=GREY, align=PP_ALIGN.CENTER)
connect(10.46, 2.5, 10.62, 2.5, color=BORDER)                       # azure -> external
textbox(10.0, 2.16, 1.6, 0.2, "via egress proxy", size=6, color=GREY, align=PP_ALIGN.CENTER)

# ---------- footer narrative ----------
fb = rect(0.35, 7.0, 12.63, 0.42, WHITE, line=BORDER)
write(fb, "Shell user to DNS (vendorpulse.it.shell.com) to Front Door (WAF) to App Service (private) to PostgreSQL via Private Link.   "
          "Graph and Foundry via outbound HTTPS (egress proxy).   Managed Identity to Key Vault / DB, app-only certificate to Graph.",
      tcolor=GREY, tsize=7.8, bold=False)

OUT = "VendorPulse_Deployment_Architecture.pptx"; prs.save(OUT); print("wrote", OUT)
