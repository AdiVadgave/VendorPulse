"""Build VendorPulse_Offering_Deck.pptx — a fresh-theme offering deck.

Storytelling flow (12 slides):
  Opening, Business Challenges, Why this Solution, What is this Solution,
  Solution Overview, Reference Architecture, Business KPIs / Outcome Metrics,
  Benefits, Success Stories, Why Zensar, Next Steps, Thank You.

Theme: modern deep-navy + teal/cyan accent (NOT the Zensar red template).
Content is grounded in the actual VendorPulse codebase analysis.

Usage: python _build_offering_deck_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = "VendorPulse_Offering_Deck.pptx"

# ---- Theme palette (navy + teal) ----
NAVY = RGBColor(0x0B, 0x1F, 0x3A)      # deep navy background
NAVY2 = RGBColor(0x12, 0x2B, 0x4E)     # panel navy
INK = RGBColor(0x1B, 0x2A, 0x41)       # body ink
TEAL = RGBColor(0x14, 0xB8, 0xA6)      # primary accent
CYAN = RGBColor(0x2B, 0xA8, 0xE0)      # secondary accent
AMBER = RGBColor(0xF4, 0xB7, 0x40)     # highlight
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF7, 0xFA)     # light slide bg
CARD = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5B, 0x6B, 0x7F)
LINE = RGBColor(0xD8, 0xDF, 0xE8)
CHIP = RGBColor(0xE7, 0xF6, 0xF4)      # teal-tint chip

HFONT = "Segoe UI Semibold"
TFONT = "Segoe UI"
BFONT = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def _set_bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return r


def _rect(slide, shape, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, *rest) in para:
            italic = rest[0] if rest else False
            run = p.add_run(); run.text = txt
            run.font.size = Pt(size); run.font.color.rgb = color
            run.font.bold = bold; run.font.name = font; run.font.italic = italic
    return tb


def _bullets(slide, l, t, w, h, items, size=13, color=INK, gap=6, marker="–",
             marker_color=TEAL, bold_lead=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.04
        m = p.add_run(); m.text = marker + "  "
        m.font.size = Pt(size); m.font.color.rgb = marker_color; m.font.bold = True; m.font.name = TFONT
        # support **bold lead** before a colon
        if bold_lead and "—" in item:
            lead, rest = item.split("—", 1)
            r1 = p.add_run(); r1.text = lead.strip() + " — "
            r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.bold = True; r1.font.name = TFONT
            r2 = p.add_run(); r2.text = rest.strip()
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.bold = False; r2.font.name = TFONT
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = TFONT
    return tb


def _header(slide, kicker, title):
    # light slide background
    _set_bg(slide, LIGHT)
    # top accent bar
    _rect(slide, MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(0.14), fill=TEAL)
    # kicker + title
    _text(slide, Inches(0.6), Inches(0.42), Inches(11.5), Inches(0.35),
          [[(kicker.upper(), 11, TEAL, True, HFONT)]])
    _text(slide, Inches(0.6), Inches(0.72), Inches(12.1), Inches(0.8),
          [[(title, 28, NAVY, True, HFONT)]])
    # underline
    _rect(slide, MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.42), Inches(1.1), Pt(3), fill=AMBER)


def _footer(slide, idx):
    _text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
          [[("VendorPulse  ·  Zensar Offering Collateral  ·  Confidential", 8, MUTED, False, TFONT)]])
    _text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.3),
          [[(str(idx), 8, MUTED, False, TFONT)]], align=PP_ALIGN.RIGHT)


def _chip(slide, l, t, text, fill=CHIP, color=NAVY, w=Inches(2.2)):
    s = _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.42), fill=fill)
    s.adjustments[0] = 0.5
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = color; r.font.name = TFONT
    return s


def _card(slide, l, t, w, h, title, body, accent=TEAL, title_size=14, body_size=11,
          title_color=NAVY, body_color=MUTED):
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill=CARD, line=LINE, line_w=1.0).adjustments[0] = 0.06
    # accent tab
    _rect(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(0.12), h, fill=accent).adjustments[0] = 0.5
    tb = slide.shapes.add_textbox(l + Inches(0.28), t + Inches(0.16), w - Inches(0.45), h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.space_after = Pt(4)
    r = p.add_run(); r.text = title; r.font.size = Pt(title_size); r.font.bold = True; r.font.color.rgb = title_color; r.font.name = HFONT
    if body:
        p2 = tf.add_paragraph(); p2.line_spacing = 1.04
        r2 = p2.add_run(); r2.text = body; r2.font.size = Pt(body_size); r2.font.color.rgb = body_color; r2.font.name = TFONT


def build():
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    # ===================== 1. OPENING =====================
    s = slide()
    _set_bg(s, NAVY)
    # accent geometry
    _rect(s, MSO_SHAPE.RECTANGLE, 0, Inches(6.9), EMU_W, Inches(0.6), fill=NAVY2)
    _rect(s, MSO_SHAPE.RECTANGLE, 0, Inches(2.55), Inches(1.5), Pt(4), fill=TEAL)
    _text(s, Inches(0.9), Inches(1.7), Inches(11), Inches(0.5),
          [[("ZENSAR  ·  AI OFFERING", 13, TEAL, True, HFONT)]])
    _text(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(1.6),
          [[("VendorPulse", 54, WHITE, True, HFONT)]])
    _text(s, Inches(0.9), Inches(3.95), Inches(11.4), Inches(0.9),
          [[("AI-assisted vendor governance for EGB & QBR cycles", 22, RGBColor(0xCF,0xDB,0xE8), False, TFONT)]])
    _text(s, Inches(0.9), Inches(4.95), Inches(11.4), Inches(0.5),
          [[("Deterministic-first.  AI-second.  Human-approved.", 15, AMBER, True, TFONT)]])
    _text(s, Inches(0.9), Inches(7.02), Inches(11.4), Inches(0.4),
          [[("Offering collateral  ·  MV BOT (Mobility VMO automation)  ·  Confidential", 10, RGBColor(0x9F,0xB0,0xC2), False, TFONT)]])

    # ===================== 2. BUSINESS CHALLENGES =====================
    s = slide()
    _header(s, "The problem", "Business Challenges in Vendor Governance")
    _text(s, Inches(0.6), Inches(1.62), Inches(12), Inches(0.4),
          [[("Every EGB / QBR cycle repeats the same manual, high-stakes coordination — quarter after quarter.", 13, MUTED, False, TFONT)]])
    cw, ch = Inches(5.95), Inches(1.85)
    gx, gy = Inches(0.6), Inches(2.25)
    gap = Inches(0.25)
    chal = [
        ("Manual, multi-party scheduling", "Coordinating executives, sponsors and vendors across calendars every quarter consumes coordinator time.", TEAL),
        ("Slow, inconsistent scorecards", "Collecting and validating vendor + internal scores is manual, error-prone and often late.", CYAN),
        ("Manual prep & lost memory", "Briefs, alignment notes and minutes are hand-built; institutional insight is lost between cycles.", AMBER),
        ("Governance & audit pressure", "Human oversight, traceability and auditability are now mandatory (Shell IRM 3.492 / EU AI Act).", RGBColor(0xE0,0x6A,0x6A)),
    ]
    for i, (ti, bo, ac) in enumerate(chal):
        col = i % 2; row = i // 2
        l = gx + col * (cw + gap); t = gy + row * (ch + gap)
        _card(s, l, t, cw, ch, ti, bo, accent=ac, title_size=16, body_size=12)
    _footer(s, 2)

    # ===================== 3. WHY THIS SOLUTION =====================
    s = slide()
    _header(s, "Why this solution", "Governance you can automate — without losing control")
    _bullets(s, Inches(0.6), Inches(1.85), Inches(7.0), Inches(4.6), [
        "Governance is recurring, structured and high-stakes — ideal for automation, but it cannot be left to a black-box AI.",
        "VendorPulse automates the busywork while keeping every decision in deterministic, auditable code.",
        "AI is confined to drafting human-readable text — it never computes scores, decides outcomes, or sends anything.",
        "A human approves every outbound action, and every run is logged — so speed never costs you control.",
        "The result: faster cycles, consistent quality, and an audit trail that satisfies governance and compliance.",
    ], size=15, gap=14, bold_lead=False)
    # callout panel
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.0), Inches(4.7), Inches(3.6), fill=NAVY).adjustments[0] = 0.05
    _text(s, Inches(8.35), Inches(2.5), Inches(4.0), Inches(2.8), [
        [("THE ONE-LINE PITCH", 11, TEAL, True, HFONT)],
        [("“The governance copilot where the AI never decides and never sends —", 17, WHITE, True, HFONT)],
        [("a human approves every action, and every step is auditable.”", 17, WHITE, True, HFONT)],
    ], space_after=8, line_spacing=1.08)
    _footer(s, 3)

    # ===================== 4. WHAT IS THIS SOLUTION =====================
    s = slide()
    _header(s, "What it is", "A deterministic-first vendor-governance platform")
    _bullets(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(3.0), [
        "An end-to-end platform that orchestrates the full EGB / QBR cycle — from attendee refresh to archived minutes.",
        "A forward-only 12-state workflow engine spanning 6 functional modules (A–F), enforced in code.",
        "6 specialized agents that draft text; all scoring, ranking and transitions are deterministic logic.",
        "A human-approval gate on every outbound action — invites, briefs and minutes only send after sign-off.",
        "Runs with or without the LLM (a config toggle) — the deterministic core always works.",
    ], size=14.5, gap=11)
    # by-the-numbers strip
    stats = [("12", "workflow states"), ("6", "modules"), ("6", "agents"), ("16", "scorecard KPIs"), ("4", "KPI categories")]
    n = len(stats); sw = Inches(2.32); sgap = Inches(0.18); total = n*sw + (n-1)*sgap
    sx = (EMU_W - total) / 2; sy = Inches(5.3)
    for i, (num, lab) in enumerate(stats):
        l = sx + i*(sw+sgap)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, sy, sw, Inches(1.25), fill=NAVY).adjustments[0]=0.08
        _text(s, l, sy+Inches(0.18), sw, Inches(0.6), [[(num, 30, TEAL, True, HFONT)]], align=PP_ALIGN.CENTER)
        _text(s, l, sy+Inches(0.82), sw, Inches(0.35), [[(lab, 11, RGBColor(0xCF,0xDB,0xE8), False, TFONT)]], align=PP_ALIGN.CENTER)
    _footer(s, 4)

    # ===================== 5. SOLUTION OVERVIEW (6 modules) =====================
    s = slide()
    _header(s, "How it works", "Solution Overview — six modules, one workflow")
    mods = [
        ("A · Scheduling", "Attendee refresh, availability, deterministic slot ranking, invite + Teams meeting on approval.", TEAL),
        ("B · Scorecard", "Form collection, deterministic validation (1–5), outlier detection (z-score 1.5 SD), compile.", CYAN),
        ("C · Alignment", "Score-diff & internal-vs-vendor spread detection, “what changed” summary, action extraction.", TEAL),
        ("D · Vendor Prep", "Grounded vendor brief + 3 pushback response options — every draft human-approved.", CYAN),
        ("E · Meeting", "Transcript parsed into 5 note types, minutes generation, action register.", TEAL),
        ("F · Analytics", "Cross-cycle trends, recurring-issue detection, leadership brief card.", CYAN),
    ]
    cw, ch = Inches(3.95), Inches(1.95); gx, gy = Inches(0.6), Inches(1.85); gpx, gpy = Inches(0.18), Inches(0.2)
    for i, (ti, bo, ac) in enumerate(mods):
        col = i % 3; row = i // 3
        l = gx + col*(cw+gpx); t = gy + row*(ch+gpy)
        _card(s, l, t, cw, ch, ti, bo, accent=ac, title_size=15, body_size=11)
    _text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.4),
          [[("Deterministic logic in every module — the AI only drafts the words; the numbers and decisions are code.", 12, NAVY, True, TFONT)]])
    _footer(s, 5)

    # ===================== 6. REFERENCE ARCHITECTURE =====================
    s = slide()
    _header(s, "Architecture", "Reference Architecture — single-tenant on Azure")
    zones = [
        ("1 · CLIENT", ["Browser (Edge/any)", "Entra SSO login", "Coordinator · Sponsor · Viewer"], CYAN),
        ("2 · EDGE", ["App Gateway + WAF", "TLS · OWASP", "Origin-lock"], TEAL),
        ("3 · APPLICATION", ["FastAPI + React 19 SPA", "Entra OIDC + RBAC", "12-state WorkflowEngine", "Approval Gate (HITL)", "6 agents + services"], NAVY),
        ("4 · DATA & AI", ["PostgreSQL · Key Vault", "Blob · App Insights (OTel)", "Azure AI Foundry GPT-4o", "(in Shell tenant)"], TEAL),
        ("5 · EXTERNAL", ["via egress proxy:", "Microsoft Graph", "(mail · calendar · Teams)", "Entra ID"], CYAN),
    ]
    zw = Inches(2.42); zgap = Inches(0.07); zx = Inches(0.55); zy = Inches(1.95); zh = Inches(3.7)
    for i, (ti, items, ac) in enumerate(zones):
        l = zx + i*(zw+zgap)
        hdr_fill = NAVY if ti.startswith("3") else NAVY2
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, zy, zw, zh, fill=RGBColor(0xFF,0xFF,0xFF), line=LINE, line_w=1.0).adjustments[0]=0.04
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, zy, zw, Inches(0.5), fill=ac).adjustments[0]=0.06
        _text(s, l, zy+Inches(0.08), zw, Inches(0.4), [[(ti, 12, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER)
        _bullets(s, l+Inches(0.18), zy+Inches(0.66), zw-Inches(0.32), zh-Inches(0.8), items, size=10.5, gap=4, marker="•", marker_color=ac, color=INK)
        if i < len(zones)-1:
            _text(s, l+zw-Inches(0.02), zy+Inches(1.6), Inches(0.18), Inches(0.4), [[("›", 20, MUTED, True, HFONT)]], align=PP_ALIGN.CENTER)
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.95), Inches(12.25), Inches(0.7), fill=CHIP).adjustments[0]=0.3
    _text(s, Inches(0.8), Inches(6.06), Inches(11.8), Inches(0.5),
          [[("AI runs in-tenant; all internal traffic stays private; the only thing leaving the network is Microsoft 365 — through one controlled, logged egress.", 12, NAVY, True, TFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, 6)

    # ===================== 7. BUSINESS KPIs / OUTCOME METRICS =====================
    s = slide()
    _header(s, "Value", "Business KPIs & Outcome Metrics")
    kpis = [
        ("~50%", "less coordinator effort per cycle (scheduling, scorecards, minutes automated)", TEAL),
        ("Weeks → days", "for scheduling and scorecard compilation", CYAN),
        ("16 KPIs", "validated deterministically every cycle, across 4 categories", TEAL),
        ("100%", "of AI outputs human-approved and logged — full audit coverage", AMBER),
        ("1.5 SD", "automatic outlier detection on every scorecard", CYAN),
        ("Cross-cycle", "recurring-issue detection and carry-forward actions", TEAL),
    ]
    cw, ch = Inches(3.95), Inches(1.85); gx, gy = Inches(0.6), Inches(1.9); gpx, gpy = Inches(0.18), Inches(0.2)
    for i, (num, lab, ac) in enumerate(kpis):
        col = i % 3; row = i // 3
        l = gx + col*(cw+gpx); t = gy + row*(ch+gpy)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, cw, ch, fill=CARD, line=LINE, line_w=1.0).adjustments[0]=0.07
        _text(s, l+Inches(0.25), t+Inches(0.22), cw-Inches(0.4), Inches(0.7), [[(num, 26, ac, True, HFONT)]])
        _text(s, l+Inches(0.25), t+Inches(0.95), cw-Inches(0.45), Inches(0.8), [[(lab, 11.5, MUTED, False, TFONT)]], line_spacing=1.05)
    _text(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.4),
          [[("Figures are target / illustrative for an internal governance workload; confirmed against pilot usage.", 10.5, MUTED, False, TFONT, True)]])
    _footer(s, 7)

    # ===================== 8. BENEFITS =====================
    s = slide()
    _header(s, "Why it matters", "Benefits")
    bens = [
        ("Speed", "Automates the chase — scheduling, reminders, compilation and minutes — freeing coordinators for judgement work.", TEAL),
        ("Consistency & quality", "Deterministic validation and grounded AI drafts (no invented numbers) make every cycle uniform and reliable.", CYAN),
        ("Trust & compliance", "Human-approval gate, full audit trail and in-tenant AI map directly to IRM 3.492 and the EU AI Act.", AMBER),
        ("Continuity", "Cross-cycle memory surfaces trends, recurring issues and carry-forward actions — each cycle builds on the last.", TEAL),
        ("Low risk", "Runs with the AI switched off; the deterministic core always works — no black-box dependency.", CYAN),
        ("Microsoft-native", "Built on Entra, Graph, Foundry and Azure — nothing novel leaves the enterprise boundary.", RGBColor(0xE0,0x6A,0x6A)),
    ]
    cw, ch = Inches(3.95), Inches(1.95); gx, gy = Inches(0.6), Inches(1.85); gpx, gpy = Inches(0.18), Inches(0.2)
    for i, (ti, bo, ac) in enumerate(bens):
        col = i % 3; row = i // 3
        l = gx + col*(cw+gpx); t = gy + row*(ch+gpy)
        _card(s, l, t, cw, ch, ti, bo, accent=ac, title_size=15, body_size=11)
    _footer(s, 8)

    # ===================== 9. SUCCESS STORIES =====================
    s = slide()
    _header(s, "Proof", "Success Stories — working MVP & pilot scenarios")
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.65), Inches(12.1), Inches(0.55), fill=CHIP).adjustments[0]=0.4
    _text(s, Inches(0.85), Inches(1.74), Inches(11.6), Inches(0.4),
          [[("Illustrative scenarios from the working full-stack MVP (3 vendors, multiple seeded cycles). Client references to be added.", 11.5, NAVY, True, TFONT)]], anchor=MSO_ANCHOR.MIDDLE)
    stories = [
        ("Scheduling, automated", "Meridian IT (Managed Services)", "Attendees refreshed, availability polled and slots deterministically ranked; Teams meeting created on approval — coordinator effort cut sharply.", TEAL),
        ("Scorecards on track", "Novatech (IT Infrastructure)", "Scorecard dispatched to stakeholders, validated deterministically and compiled — no manual chasing or spreadsheet errors.", CYAN),
        ("Full cycle, audited", "Archived QBR cycle", "End-to-end run through all 12 states with minutes and an immutable agent-run audit trail — every AI draft human-approved.", AMBER),
    ]
    cw, ch = Inches(3.95), Inches(3.4); gx, gy = Inches(0.6), Inches(2.45); gpx = Inches(0.18)
    for i, (ti, sub, bo, ac) in enumerate(stories):
        l = gx + i*(cw+gpx)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, gy, cw, ch, fill=CARD, line=LINE, line_w=1.0).adjustments[0]=0.05
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, gy, cw, Inches(0.12), fill=ac).adjustments[0]=0.5
        _text(s, l+Inches(0.28), gy+Inches(0.35), cw-Inches(0.5), Inches(0.5), [[(ti, 16, NAVY, True, HFONT)]])
        _text(s, l+Inches(0.28), gy+Inches(0.85), cw-Inches(0.5), Inches(0.4), [[(sub, 11.5, ac, True, TFONT)]])
        _text(s, l+Inches(0.28), gy+Inches(1.4), cw-Inches(0.55), Inches(1.8), [[(bo, 12, MUTED, False, TFONT)]], line_spacing=1.08)
    _footer(s, 9)

    # ===================== 10. WHY ZENSAR =====================
    s = slide()
    _header(s, "The partner", "Why Zensar")
    pillars = [
        ("AI-first, governance-led", "We build AI for regulated enterprises — human-in-the-loop and auditability by design, not bolted on.", TEAL),
        ("Microsoft Azure & Foundry depth", "Hands-on with Entra, Graph, Foundry, App Gateway and Azure data services across the stack.", CYAN),
        ("A working MVP, not slideware", "Full-stack build already running: FastAPI backend, React 19 SPA, 6 agents, a 12-state engine.", AMBER),
        ("Built for Shell-grade compliance", "Designed against IRM 3.492 and the EU AI Act — deterministic core, in-tenant AI, full audit.", TEAL),
        ("Ready-to-deploy accelerator", "Config-driven and modular — adapt to your vendors, KPIs and cadence quickly.", CYAN),
        ("De-risked delivery", "Deterministic fallback means the system works with or without the LLM from day one.", RGBColor(0xE0,0x6A,0x6A)),
    ]
    cw, ch = Inches(3.95), Inches(1.95); gx, gy = Inches(0.6), Inches(1.85); gpx, gpy = Inches(0.18), Inches(0.2)
    for i, (ti, bo, ac) in enumerate(pillars):
        col = i % 3; row = i // 3
        l = gx + col*(cw+gpx); t = gy + row*(ch+gpy)
        _card(s, l, t, cw, ch, ti, bo, accent=ac, title_size=14, body_size=11)
    _footer(s, 10)

    # ===================== 11. NEXT STEPS =====================
    s = slide()
    _header(s, "Let's go", "Next Steps")
    steps = [
        ("01", "Discovery workshop", "1–2 weeks", "Confirm scope, integrations, KPIs and compliance prerequisites; agree the roadmap.", TEAL),
        ("02", "Foundations & pilot", "3–4 weeks", "Stand up Azure infra, Entra/Graph, migrate to PostgreSQL, port agents, run one live cycle.", CYAN),
        ("03", "UAT & go-live", "2 weeks", "VMO UAT, IT-security sign-off, production cutover, coordinator training and handover.", AMBER),
    ]
    cw, ch = Inches(3.95), Inches(3.0); gx, gy = Inches(0.6), Inches(1.95); gpx = Inches(0.18)
    for i, (no, ti, dur, bo, ac) in enumerate(steps):
        l = gx + i*(cw+gpx)
        _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, l, gy, cw, ch, fill=CARD, line=LINE, line_w=1.0).adjustments[0]=0.05
        _text(s, l+Inches(0.28), gy+Inches(0.28), Inches(2), Inches(0.8), [[(no, 34, ac, True, HFONT)]])
        _text(s, l+Inches(0.28), gy+Inches(1.05), cw-Inches(0.5), Inches(0.45), [[(ti, 16, NAVY, True, HFONT)]])
        _chip(s, l+Inches(0.28), gy+Inches(1.55), dur, fill=CHIP, color=NAVY, w=Inches(1.6))
        _text(s, l+Inches(0.28), gy+Inches(2.15), cw-Inches(0.55), Inches(0.8), [[(bo, 11.5, MUTED, False, TFONT)]], line_spacing=1.06)
    _rect(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.95), fill=NAVY).adjustments[0]=0.12
    _text(s, Inches(0.9), Inches(5.32), Inches(11.5), Inches(0.7), [
        [("In parallel (start now — external lead time): ", 13, AMBER, True, TFONT), ("Shell compliance — AI Registry, IRM IAQ, Shell.AI + TRB approval — must clear before production.", 13, WHITE, False, TFONT)],
    ], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    _text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
          [[("Let’s run a discovery workshop and co-build the roadmap.", 14, NAVY, True, HFONT)]])
    _footer(s, 11)

    # ===================== 12. THANK YOU =====================
    s = slide()
    _set_bg(s, NAVY)
    _rect(s, MSO_SHAPE.RECTANGLE, 0, Inches(3.0), Inches(1.5), Pt(4), fill=TEAL)
    _text(s, Inches(0.9), Inches(2.6), Inches(11), Inches(1.2), [[("Thank You", 50, WHITE, True, HFONT)]])
    _text(s, Inches(0.92), Inches(3.9), Inches(11), Inches(0.6),
          [[("VendorPulse — AI-assisted, human-approved vendor governance", 18, RGBColor(0xCF,0xDB,0xE8), False, TFONT)]])
    _text(s, Inches(0.92), Inches(4.7), Inches(11), Inches(0.5),
          [[("Zensar  ·  www.zensar.com  ·  Confidential", 13, TEAL, True, TFONT)]])

    prs.save(OUT)
    print("Wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    build()
