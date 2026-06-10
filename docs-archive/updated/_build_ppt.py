# -*- coding: utf-8 -*-
"""Build ZenVendorPulse 3-slide deck: Architecture, Flow, Tech Stack."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette (pleasant teal / slate) ----
TEAL      = RGBColor(0x0F, 0x76, 0x6E)   # primary
TEAL_DK   = RGBColor(0x0B, 0x52, 0x4B)
TEAL_LT   = RGBColor(0xCC, 0xE9, 0xE6)
SLATE     = RGBColor(0x1E, 0x29, 0x3B)   # heading text
SLATE_MD  = RGBColor(0x47, 0x55, 0x69)   # body text
MIST      = RGBColor(0xF1, 0xF5, 0xF9)   # card bg
MIST2     = RGBColor(0xE7, 0xEE, 0xF3)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # accent
CORAL     = RGBColor(0xE1, 0x6A, 0x54)
INK_BG    = RGBColor(0xF7, 0xFA, 0xFA)   # slide bg
LINE      = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = INK_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def no_shadow(sh):
    sh.shadow.inherit = False

def soft_shadow(sh):
    """Apply a subtle drop shadow."""
    spPr = sh._element.spPr
    el = spPr.find(qn('a:effectLst'))
    if el is None:
        el = spPr.makeelement(qn('a:effectLst'), {})
        spPr.append(el)
    sh_el = el.makeelement(qn('a:outerShdw'),
        {'blurRad':'70000','dist':'30000','dir':'5400000','rotWithShape':'0'})
    clr = sh_el.makeelement(qn('a:srgbClr'), {'val':'1E293B'})
    alpha = clr.makeelement(qn('a:alpha'), {'val':'18000'})
    clr.append(alpha); sh_el.append(clr); el.append(sh_el)

def txt(s, l, t, w, h, text, size=14, color=SLATE, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, line_spacing=1.0):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.name = font
        f.bold = bold; f.italic = italic; f.color.rgb = color
    return tb

def box(s, l, t, w, h, fill, line_clr=None, line_w=1.0, radius=True, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_clr is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_clr; shp.line.width = Pt(line_w)
    no_shadow(shp)
    if shadow: soft_shadow(shp)
    # tighten corner radius
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp

def label_box(s, l, t, w, h, title, body, fill, title_clr, body_clr=None,
              tsize=13, bsize=10.5, shadow=True, line_clr=None):
    b = box(s, l, t, w, h, fill, line_clr=line_clr, shadow=shadow)
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(5); tf.margin_bottom = Pt(5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = title_clr
    if body:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.line_spacing = 1.0
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(bsize); r2.font.name = FONT
        r2.font.color.rgb = body_clr or SLATE_MD
    return b

def connector(s, x1, y1, x2, y2, color=TEAL, w=2.0, dashed=False, arrow=True):
    cn = s.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    no_shadow(cn)
    lnEl = cn.line._get_or_add_ln()
    if arrow:
        end = lnEl.makeelement(qn('a:tailEnd'),
            {'type':'triangle','w':'med','len':'med'})
        lnEl.append(end)
    if dashed:
        d = lnEl.makeelement(qn('a:prstDash'), {'val':'dash'})
        lnEl.append(d)
    return cn

def header(s, kicker, title, page):
    # top accent bar
    bar = box(s, 0, 0, SW, Inches(0.16), TEAL, radius=False)
    # kicker
    txt(s, Inches(0.6), Inches(0.34), Inches(9), Inches(0.3),
        kicker.upper(), size=11.5, color=TEAL, bold=True)
    txt(s, Inches(0.6), Inches(0.6), Inches(11.4), Inches(0.7),
        title, size=27, color=SLATE, bold=True)
    # brand chip top-right
    chip = box(s, Inches(11.0), Inches(0.34), Inches(1.75), Inches(0.42), TEAL_LT,
               shadow=False)
    txt(s, Inches(11.0), Inches(0.34), Inches(1.75), Inches(0.42),
        "ZenVendorPulse", size=10.5, color=TEAL_DK, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # page no
    txt(s, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
        str(page), size=10, color=LINE, align=PP_ALIGN.RIGHT)
    txt(s, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        "Shell  ·  Vendor Governance Automation (CR-9 EGB)", size=9.5, color=LINE)

# ============================================================
# SLIDE 1 — INTRODUCTION / COVER
# ============================================================
s = slide()
# hero band
hero = box(s, 0, 0, SW, Inches(2.85), TEAL, radius=False)
no_shadow(hero)
# accent stripe at bottom of hero
box(s, 0, Inches(2.85), SW, Inches(0.10), AMBER, radius=False)
# brand chip
chip = box(s, Inches(0.7), Inches(0.55), Inches(2.2), Inches(0.5), TEAL_DK, shadow=False)
txt(s, Inches(0.7), Inches(0.55), Inches(2.2), Inches(0.5), "ZENSAR  ·  CR-9 EGB",
    size=11, color=TEAL_LT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.95),
    "ZenVendorPulse", size=46, color=WHITE, bold=True)
txt(s, Inches(0.73), Inches(2.12), Inches(12), Inches(0.6),
    "Agentic AI for Shell Vendor Governance  ·  EGB / QBR Automation",
    size=18, color=TEAL_LT, bold=False)

# challenge / solution two-column
cy = Inches(3.35)
ch_ = Inches(1.65)
chal = box(s, Inches(0.7), cy, Inches(5.9), ch_, WHITE, line_clr=LINE, shadow=True)
txt(s, Inches(0.95), cy+Inches(0.14), Inches(5.5), Inches(0.35),
    "THE CHALLENGE", size=12, color=CORAL, bold=True)
txt(s, Inches(0.95), cy+Inches(0.52), Inches(5.45), Inches(1.05),
    "Vendor governance runs manually — availability chased over email, scorecards "
    "scattered across inboxes, alignment calls re-hashing recurring issues, and "
    "inconsistent minutes quarter on quarter.",
    size=11.5, color=SLATE_MD, line_spacing=1.05)

sol = box(s, Inches(6.75), cy, Inches(5.9), ch_, WHITE, line_clr=TEAL, line_w=1.5, shadow=True)
txt(s, Inches(7.0), cy+Inches(0.14), Inches(5.5), Inches(0.35),
    "THE SOLUTION", size=12, color=TEAL, bold=True)
txt(s, Inches(7.0), cy+Inches(0.52), Inches(5.45), Inches(1.05),
    "An agentic platform that orchestrates the full cycle end-to-end while keeping "
    "every outbound action under human approval — deployed single-tenant in Shell's "
    "Azure tenant, integrated with Outlook + Teams.",
    size=11.5, color=SLATE_MD, line_spacing=1.05)

# metric chips row
metrics = [
    ("6", "Governance\ncapabilities (A–F)"),
    ("12", "State forward-only\nworkflow machine"),
    ("4 wks", "Sprint-paced\ndelivery"),
    ("1", "Single-tenant\nAzure / M365"),
]
mw = Inches(2.86); mgap = Inches(0.30); mx = Inches(0.7); my = Inches(5.35)
for i,(num,lab) in enumerate(metrics):
    x = mx + i*(mw+mgap)
    mb = box(s, x, my, mw, Inches(1.15), MIST, line_clr=TEAL_LT, shadow=True)
    txt(s, x, my+Inches(0.12), mw, Inches(0.55), num, size=30, color=TEAL, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, x, my+Inches(0.68), mw, Inches(0.42), lab, size=10.5, color=SLATE_MD,
        align=PP_ALIGN.CENTER, line_spacing=0.95)

txt(s, Inches(0.7), Inches(6.95), Inches(11), Inches(0.3),
    "Shell  ·  Vendor Governance Automation (CR-9 EGB)", size=9.5, color=LINE)
txt(s, Inches(12.4), Inches(6.95), Inches(0.7), Inches(0.3), "1", size=10,
    color=LINE, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2 — SOLUTION ARCHITECTURE
# ============================================================
s = slide()
header(s, "Solution Architecture", "Single-tenant on Shell's Azure & Microsoft 365", 2)

# ---- Shell Users (external entry, above the Azure boundary) ----
label_box(s, Inches(0.75), Inches(1.52), Inches(2.6), Inches(0.46),
          "Shell Users  ·  Browser", "Entra SSO (OIDC) · RBAC", WHITE, SLATE,
          line_clr=LINE, tsize=11.5, bsize=9, shadow=True)

top = Inches(2.12)
CB = Inches(4.32)   # container height
# ---- Outer Azure subscription container ----
sub = box(s, Inches(0.55), top, Inches(8.55), CB, WHITE,
          line_clr=TEAL, line_w=1.5, shadow=True)
txt(s, Inches(0.75), top+Inches(0.07), Inches(8), Inches(0.3),
    "SHELL AZURE SUBSCRIPTION  ·  West Europe  ·  VNet + Private Link",
    size=10, color=TEAL_DK, bold=True)

# Front Door
fd = label_box(s, Inches(0.78), top+Inches(0.52), Inches(1.95), Inches(0.92),
          "Azure Front Door", "WAF + TLS\nOrigin lock", TEAL, WHITE, body_clr=TEAL_LT,
          tsize=12, bsize=9.5)

# App Service container
appc = box(s, Inches(3.0), top+Inches(0.46), Inches(3.15), Inches(2.55), MIST,
           line_clr=TEAL, line_w=1.2, shadow=True)
txt(s, Inches(3.0), top+Inches(0.52), Inches(3.15), Inches(0.3),
    "Azure App Service · Linux containers", size=10, color=TEAL_DK, bold=True,
    align=PP_ALIGN.CENTER)
label_box(s, Inches(3.2), top+Inches(0.96), Inches(2.75), Inches(0.7),
          "Frontend — React SPA", "TypeScript · nginx", WHITE, SLATE, line_clr=LINE,
          tsize=11.5, bsize=9.5, shadow=False)
label_box(s, Inches(3.2), top+Inches(1.78), Inches(2.75), Inches(1.02),
          "FastAPI Backend", "Python 3.11 · 6 AI agents\nWorkflow engine + approvals",
          TEAL_LT, TEAL_DK, body_clr=SLATE_MD, tsize=11.5, bsize=9, shadow=False)

# Managed identity note
txt(s, Inches(3.0), top+Inches(3.04), Inches(3.15), Inches(0.28),
    "↓ Managed Identity — no secrets in code", size=8.5, color=SLATE_MD,
    align=PP_ALIGN.CENTER, italic=True)

# data layer trio
dy = top+Inches(3.35)
label_box(s, Inches(0.75), dy, Inches(2.5), Inches(0.86),
          "Azure Key Vault", "LLM key · Graph cert", WHITE, SLATE,
          line_clr=LINE, tsize=11, bsize=9, shadow=False)
label_box(s, Inches(3.38), dy, Inches(2.5), Inches(0.86),
          "PostgreSQL", "Flexible Server · HA\nPrivate Link", WHITE, SLATE,
          line_clr=LINE, tsize=11, bsize=8.5, shadow=False)
label_box(s, Inches(6.0), dy, Inches(2.9), Inches(0.86),
          "App Insights", "Log Analytics\nAudit · observability", WHITE, SLATE,
          line_clr=LINE, tsize=11, bsize=8.5, shadow=False)

# ---- External services column ----
ex = box(s, Inches(9.45), top, Inches(3.35), CB, MIST2,
         line_clr=LINE, line_w=1.0, shadow=True)
txt(s, Inches(9.6), top+Inches(0.07), Inches(3), Inches(0.3),
    "EXTERNAL SERVICES  ·  outbound HTTPS", size=9.5, color=SLATE_MD, bold=True)
label_box(s, Inches(9.62), top+Inches(0.5), Inches(3.0), Inches(0.86),
          "Microsoft Entra ID", "Shell SSO (OIDC) · RBAC", WHITE, SLATE,
          line_clr=TEAL, tsize=12, bsize=9.5)
label_box(s, Inches(9.62), top+Inches(1.5), Inches(3.0), Inches(1.12),
          "Microsoft Graph API", "Mail.Send · Calendars.RW\nOnlineMeetings · Teams\nApp-only certificate auth",
          WHITE, SLATE, line_clr=TEAL, tsize=12, bsize=9)
label_box(s, Inches(9.62), top+Inches(2.76), Inches(3.0), Inches(1.42),
          "LLM Provider", "Anthropic Claude (Sonnet 4.6)\nor Azure OpenAI · switchable\nHuman-readable text only",
          TEAL, WHITE, body_clr=TEAL_LT, tsize=12, bsize=9)

# connectors
connector(s, Inches(2.05), Inches(1.98), Inches(1.75), top+Inches(0.52), TEAL, 2.0)   # user->FD
connector(s, Inches(2.73), top+Inches(0.98), Inches(3.0), top+Inches(1.31), TEAL, 2.0) # FD->app
connector(s, Inches(4.57), top+Inches(3.01), Inches(4.57), dy, LINE, 1.5)             # app->data
connector(s, Inches(6.15), top+Inches(1.7), Inches(9.45), top+Inches(0.93), TEAL, 1.75, dashed=True)  # ->entra
connector(s, Inches(6.15), top+Inches(2.1), Inches(9.45), top+Inches(2.06), TEAL, 1.75, dashed=True)  # ->graph
connector(s, Inches(6.15), top+Inches(2.5), Inches(9.45), top+Inches(3.47), TEAL, 1.75, dashed=True)  # ->llm

# bottom principle strip
txt(s, Inches(0.55), Inches(6.62), Inches(12.3), Inches(0.35),
    "Deterministic-first, AI-second   ·   Human-approval gate on every outbound action   ·   All secrets in Key Vault   ·   Full audit trail",
    size=10.5, color=TEAL_DK, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3 — DETAILED ARCHITECTURE (HLD)
# ============================================================
s = slide()
header(s, "Detailed Solution Architecture (HLD)", "Layered application design — frontend to data", 3)

# helper: a horizontal layer band with a left label tag + content
def layer(y, h, tag, tag_clr, content_runs, band_fill=WHITE, band_line=TEAL_LT):
    lx = Inches(0.55); lw = Inches(8.45)
    tagw = Inches(1.75)
    band = box(s, lx, y, lw, h, band_fill, line_clr=band_line, line_w=1.25, shadow=True)
    tagb = box(s, lx, y, tagw, h, tag_clr, radius=True, shadow=False)
    try: tagb.adjustments[0]=0.06
    except Exception: pass
    txt(s, lx, y, tagw, h, tag, size=11, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
    # content text frame
    cb = s.shapes.add_textbox(lx+tagw+Inches(0.12), y, lw-tagw-Inches(0.25), h)
    tf = cb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,(t,b,c,sz) in enumerate(content_runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=1.0
        r=p.add_run(); r.text=t
        r.font.size=Pt(sz); r.font.name=FONT; r.font.bold=b; r.font.color.rgb=c
    return band

y = Inches(1.62); g = Inches(0.115)
# 1. Frontend / Presentation
h1=Inches(0.82)
layer(y, h1, "PRESENTATION\nFrontend", TEAL,
      [("React 19 · Vite · TypeScript · Tailwind v4 · shadcn/ui", True, SLATE, 11),
       ("TanStack Query (server state)  +  Zustand (UI state)  ·  Axios typed client", False, SLATE_MD, 10),
       ("HttpOnly session cookie — no tokens in JS", False, SLATE_MD, 9.5)])
y += h1+g
# 2. Edge / hosting
h2=Inches(0.5)
layer(y, h2, "EDGE / HOST", TEAL_DK,
      [("Azure Front Door (WAF + TLS, origin lock)  ·  nginx static  ·  App Service Linux containers", True, SLATE, 10.5)])
y += h2+g
# 3. Middleware
h3=Inches(0.5)
layer(y, h3, "MIDDLEWARE", SLATE,
      [("FastAPI ASGI:  OIDC session validation · correlation-ID · rate limit (60/min) · logging · exception mapping", False, SLATE_MD, 10.5)])
y += h3+g
# 4. API Routes
h4=Inches(0.5)
layer(y, h4, "API ROUTES\n7 routers", TEAL,
      [("cycles · scheduling · scorecard · alignment · vendor_prep · meeting · analytics  + admin", True, SLATE, 10.5)])
y += h4+g
# 5. Orchestration
h5=Inches(0.66)
layer(y, h5, "ORCHESTRATION", AMBER,
      [("Workflow Engine — 12-state forward-only FSM (gates every endpoint)", True, SLATE, 10.5),
       ("AI Agents A–F — Claude tool-calling · structured AgentResponse contract", False, SLATE_MD, 10)])
y += h5+g
# 6. Service Layer
h6=Inches(0.66)
layer(y, h6, "SERVICE LAYER", TEAL,
      [("LLMService (provider-abstracted) · GraphService · ScorecardFormService · ValidationService", False, SLATE_MD, 10),
       ("AnalyticsService · SlotRankingService · AuditService · NotificationService", False, SLATE_MD, 10)])
y += h6+g
# 7. Data Access
h7=Inches(0.5)
layer(y, h7, "DATA ACCESS", TEAL_DK,
      [("Repositories → SQLAlchemy 2.0 (async) → asyncpg → PostgreSQL Flexible Server (15 tables, Private Link)", True, SLATE, 10.5)])

# ---- Right column: cross-cutting & external ----
rx = Inches(9.2); rw = Inches(3.6)
txt(s, rx, Inches(1.6), rw, Inches(0.3), "CROSS-CUTTING & EXTERNAL",
    size=10, color=SLATE_MD, bold=True)
cross = [
    ("Microsoft Entra ID", "SSO (OIDC) · RBAC 4 roles", TEAL),
    ("Microsoft Graph", "Mail · Calendar · Teams\napp-only certificate", TEAL),
    ("LLM Provider", "Claude / Azure OpenAI\nconfig-switchable", TEAL),
    ("Azure Key Vault", "secrets via Managed Identity", TEAL_DK),
    ("App Insights + Log Analytics", "agent_runs · external_calls\nfull audit trail", TEAL_DK),
]
cyv = Inches(1.95); cbh = Inches(0.82); cg = Inches(0.10)
for title, body, clr in cross:
    label_box(s, rx, cyv, rw, cbh, title, body, WHITE, SLATE, line_clr=clr,
              tsize=11.5, bsize=9, shadow=True)
    cyv += cbh+cg

# human-approval gate callout
gate = box(s, rx, cyv, rw, Inches(0.5), AMBER, shadow=True)
txt(s, rx, cyv, rw, Inches(0.5), "⛔  Human-approval gate — hard gate on every send",
    size=10, color=SLATE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# SLIDE 4 — FLOW DIAGRAM (6-stage workflow A-F)
# ============================================================
s = slide()
header(s, "Workflow Flow", "Six-stage vendor governance cycle (12-state machine)", 4)

stages = [
    ("A", "Meeting Scheduling", "Propose & book EGB/QBR via Graph; coordinate attendees from service mailbox"),
    ("B", "Scorecard Input", "Collect & validate 1–5 ratings: 4 categories / 16 parameters"),
    ("C", "Internal Alignment", "Consolidate scorecards; surface Shell-side discussion points"),
    ("D", "Vendor Prep", "Draft pushbacks in 3 stances — Factual · Neutral · Escalation"),
    ("E", "EGB/QBR Meeting", "Run combined meeting; capture decisions, actions & minutes"),
    ("F", "Trend Analysis", "Retain cross-cycle memory; track vendor trends & recurring themes"),
]

# two rows of 3, snaking flow
cw, ch = Inches(3.75), Inches(1.62)
gx, gy = Inches(0.30), Inches(0.55)
x0, y0 = Inches(0.62), Inches(1.85)
positions = []
for i in range(6):
    row = 0 if i < 3 else 1
    col = i % 3 if row == 0 else (2 - (i % 3))  # snake on row 2
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy + Inches(0.35))
    positions.append((x, y))

for i, (code, title, body) in enumerate(stages):
    x, y = positions[i]
    card = box(s, x, y, cw, ch, WHITE, line_clr=TEAL_LT, line_w=1.5, shadow=True)
    # accent left bar
    accent = box(s, x, y, Inches(0.14), ch, TEAL, radius=False)
    # stage badge circle
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.28), y+Inches(0.22),
                               Inches(0.62), Inches(0.62))
    badge.fill.solid(); badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background(); no_shadow(badge)
    btf = badge.text_frame; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = code
    br.font.size = Pt(22); br.font.bold = True; br.font.color.rgb = WHITE; br.font.name = FONT
    # title + body
    txt(s, x+Inches(1.02), y+Inches(0.16), cw-Inches(1.15), Inches(0.5),
        title, size=14.5, color=SLATE, bold=True)
    txt(s, x+Inches(1.02), y+Inches(0.66), cw-Inches(1.2), Inches(0.9),
        body, size=10, color=SLATE_MD, line_spacing=1.02)

# arrows: A->B->C (row1 L->R), C->D (down), D->E->F (row2 R->L since snaked)
def arrow_between(i, j, side):
    xi, yi = positions[i]; xj, yj = positions[j]
    hc = cw // 2; vc = ch // 2
    if side == 'r':  # left-to-right
        connector(s, xi+cw, yi+vc, xj, yj+vc, TEAL, 2.5)
    elif side == 'l':  # right-to-left
        connector(s, xi, yi+vc, xj+cw, yj+vc, TEAL, 2.5)
    elif side == 'd':  # downward
        connector(s, xi+hc, yi+ch, xj+hc, yj, AMBER, 2.5)

arrow_between(0,1,'r')
arrow_between(1,2,'r')
arrow_between(2,3,'d')   # C (col2,row0) down to D (col2,row1)
arrow_between(3,4,'l')
arrow_between(4,5,'l')

# loop-back F -> A (cross-cycle memory) curved hint
txt(s, Inches(0.62), Inches(6.78), Inches(12.0), Inches(0.4),
    "↻  Cross-cycle memory feeds the next governance cycle    |    Every outbound step (invites · briefs · minutes) passes a human-approval gate",
    size=10.5, color=TEAL_DK, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3 — TECHNOLOGY STACK TABLE
# ============================================================
s = slide()
header(s, "Technology Stack", "Production technology & Azure services", 5)

rows = [
    ("Frontend", "React (TypeScript)"),
    ("Backend / Agents", "Python 3.11+ · FastAPI"),
    ("Database", "PostgreSQL Flexible Server (Azure)"),
    ("Identity & Auth", "Entra ID SSO (OIDC) · app-only certificate auth"),
    ("Messaging & Calendar", "Microsoft Graph API (Mail · Calendar · Teams)"),
    ("Service Mailbox", "vendorpulse-svc@shell.com (M365 E3/E5)"),
    ("LLM Provider", "Anthropic Claude (Sonnet 4.6) or Azure OpenAI — config-switchable"),
    ("Secret Management", "Azure Key Vault"),
    ("Containerization", "Azure Container Registry + App Service (P1v3 prod / B2 nonprod)"),
    ("Observability", "Azure App Insights + Log Analytics"),
    ("CDN / WAF", "Azure Front Door with WAF"),
    ("CI/CD", "Azure DevOps or GitHub Enterprise"),
    ("Code Quality & Security", "SonarQube · TruffleHog/GitLeaks · Trivy/Snyk"),
    ("Browser Support", "Chrome, Edge (web-only · English-only)"),
]

tbl_l, tbl_t = Inches(0.62), Inches(1.62)
tbl_w = Inches(12.1)
hdr_h = Inches(0.40)
row_h = Inches(0.325)
c1_w = Inches(3.7)
c2_w = tbl_w - c1_w

# header row
box(s, tbl_l, tbl_t, c1_w, hdr_h, TEAL_DK, radius=False)
box(s, tbl_l+c1_w, tbl_t, c2_w, hdr_h, TEAL, radius=False)
txt(s, tbl_l+Inches(0.18), tbl_t, c1_w, hdr_h, "Application Component",
    size=12.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, tbl_l+c1_w+Inches(0.18), tbl_t, c2_w, hdr_h, "Technology / Service",
    size=12.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

y = tbl_t + hdr_h
for i, (comp, tech) in enumerate(rows):
    bg1 = MIST if i % 2 == 0 else WHITE
    bg2 = MIST if i % 2 == 0 else WHITE
    box(s, tbl_l, y, c1_w, row_h, bg1, radius=False)
    box(s, tbl_l+c1_w, y, c2_w, row_h, bg2, radius=False)
    txt(s, tbl_l+Inches(0.18), y, c1_w-Inches(0.2), row_h, comp,
        size=11, color=SLATE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, tbl_l+c1_w+Inches(0.18), y, c2_w-Inches(0.25), row_h, tech,
        size=11, color=SLATE_MD, anchor=MSO_ANCHOR.MIDDLE)
    y = y + row_h

# left accent border of table
box(s, tbl_l, tbl_t, Inches(0.06), y - tbl_t, AMBER, radius=False)

# footer note (environments)
txt(s, tbl_l, y + Inches(0.12), tbl_w, Inches(0.4),
    "Environments:   NonProd (B2 · standard Postgres, dev/UAT)      ·      Prod (P1v3 · Postgres HA)   —   both in Shell-designated region",
    size=10.5, color=TEAL_DK, bold=True)

out = "docs/updated/ZenVendorPulse_Architecture_Flow_TechStack.pptx"
prs.save(out)
print("Saved", out)
